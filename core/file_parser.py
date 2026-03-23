from __future__ import annotations

import logging
import re
from pathlib import Path
from urllib.parse import unquote, urlparse


logger = logging.getLogger(__name__)

SUPPORTED_EXTENSIONS = {".pdf", ".docx", ".pptx", ".txt", ".md"}
_MAX_FILE_BYTES = 12 * 1024 * 1024
_WINDOWS_PATH_PATTERN = re.compile(
    r"([a-zA-Z]:\\[^\n\r\t\"<>|]+?\.(?:pdf|docx|pptx|txt|md))",
    re.IGNORECASE,
)


def infer_file_path(value: str | None) -> Path | None:
    text = str(value or "").strip()
    if not text:
        return None
    parsed = urlparse(text)
    if parsed.scheme == "file":
        candidate = Path(unquote(parsed.path.lstrip("/")))
        return candidate if candidate.exists() and candidate.suffix.casefold() in SUPPORTED_EXTENSIONS else None

    match = _WINDOWS_PATH_PATTERN.search(text)
    if match:
        candidate = Path(match.group(1))
        if candidate.exists():
            return candidate

    direct = Path(text.strip("\"'"))
    if direct.exists() and direct.suffix.casefold() in SUPPORTED_EXTENSIONS:
        return direct
    return None


def _read_pdf(path: Path) -> str:
    try:
        import pdfplumber  # type: ignore
    except Exception:
        return ""
    parts: list[str] = []
    try:
        with pdfplumber.open(path) as pdf:
            for page in pdf.pages[:24]:
                text = (page.extract_text() or "").strip()
                if text:
                    parts.append(text)
    except Exception:
        logger.exception("Failed to read PDF text from %s", path)
        return ""
    return "\n".join(parts)


def _read_docx(path: Path) -> str:
    try:
        from docx import Document  # type: ignore
    except Exception:
        return ""
    try:
        document = Document(path)
        return "\n".join(paragraph.text.strip() for paragraph in document.paragraphs if paragraph.text.strip())
    except Exception:
        logger.exception("Failed to read DOCX text from %s", path)
        return ""


def _read_pptx(path: Path) -> str:
    try:
        from pptx import Presentation  # type: ignore
    except Exception:
        return ""
    chunks: list[str] = []
    try:
        presentation = Presentation(path)
        for slide in presentation.slides[:40]:
            for shape in slide.shapes:
                text = getattr(shape, "text", "")
                if text:
                    cleaned = str(text).strip()
                    if cleaned:
                        chunks.append(cleaned)
    except Exception:
        logger.exception("Failed to read PPTX text from %s", path)
        return ""
    return "\n".join(chunks)


def _read_text(path: Path) -> str:
    try:
        return path.read_text(encoding="utf-8", errors="ignore")
    except Exception:
        logger.exception("Failed to read text from %s", path)
        return ""


def extract_file_text(path: str | Path | None, *, max_chars: int = 8000) -> str | None:
    file_path = Path(path) if path is not None else None
    if file_path is None or not file_path.exists():
        return None
    if file_path.suffix.casefold() not in SUPPORTED_EXTENSIONS:
        return None
    try:
        if file_path.stat().st_size > _MAX_FILE_BYTES:
            return None
    except Exception:
        return None

    suffix = file_path.suffix.casefold()
    if suffix == ".pdf":
        text = _read_pdf(file_path)
    elif suffix == ".docx":
        text = _read_docx(file_path)
    elif suffix == ".pptx":
        text = _read_pptx(file_path)
    else:
        text = _read_text(file_path)

    normalized = " ".join(str(text or "").split()).strip()
    if not normalized:
        return None
    return normalized[:max_chars]
